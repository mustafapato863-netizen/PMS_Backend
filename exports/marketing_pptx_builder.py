import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# ─── COLOR PALETTE (LIGHT CORPORATE THEME) ───
BG_COLOR = RGBColor(248, 250, 252)     # Slate 50
CARD_BG = RGBColor(255, 255, 255)      # White
TEXT_DARK = RGBColor(15, 23, 42)       # Slate 900
TEXT_MUTED = RGBColor(71, 85, 105)     # Slate 600
TEXT_LIGHT = RGBColor(100, 116, 139)   # Slate 500

ACCENT_BLUE = RGBColor(2, 106, 167)
ACCENT_GREEN = RGBColor(22, 163, 74)
ACCENT_RED = RGBColor(220, 38, 38)
ACCENT_AMBER = RGBColor(217, 119, 6)
ACCENT_PURPLE = RGBColor(124, 58, 237)
ACCENT_CYAN = RGBColor(8, 145, 178)

BORDER_COLOR = RGBColor(226, 232, 240)

def set_slide_bg(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def add_kpi_card(slide, left, top, width, height, label, value, trend_text="", status_color=ACCENT_GREEN, fill_color=CARD_BG):
    card = add_rounded_card(slide, left, top, width, height, fill_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3), label, font_size=11, color=TEXT_MUTED, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.4), value, font_size=26, color=TEXT_DARK, bold=True)
    if trend_text:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.25), trend_text, font_size=10, color=status_color, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = status_color
    bar.line.fill.background()
    return card

def add_section_divider(slide, left, top, width, title, subtitle="", color=ACCENT_BLUE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    add_text_box(slide, left, top + Inches(0.1), width, Inches(0.4), title, font_size=20, color=TEXT_DARK, bold=True)
    if subtitle:
        add_text_box(slide, left, top + Inches(0.5), width, Inches(0.3), subtitle, font_size=12, color=TEXT_MUTED)

def grade_from_score(score):
    if score >= 0.95: return ('A+', ACCENT_GREEN)
    if score >= 0.90: return ('A', ACCENT_GREEN)
    if score >= 0.85: return ('B+', ACCENT_CYAN)
    if score >= 0.80: return ('B', ACCENT_BLUE)
    if score >= 0.75: return ('C+', ACCENT_AMBER)
    if score >= 0.70: return ('C', ACCENT_AMBER)
    return ('D', ACCENT_RED)


def build_marketing_pptx(period_label="June 2026") -> bytes:
    # ─── LOAD DATA ───
    df = pd.read_excel(r'D:\Trend\PMS_Trend_All.xlsx', sheet_name='Marketing')
    df['Date'] = pd.to_datetime(df['Date'])
    
    # We assume period_label is something like "June 2026"
    # To keep it generic but bound to our specific request, we just filter month 5 and 6
    may = df[df['Date'].dt.month == 5].copy()
    jun = df[df['Date'].dt.month == 6].copy()

    curr_month_df = jun
    prev_month_df = may

    perf_curr = curr_month_df.groupby(['Employee Name', 'Position']).agg({
        'Performance Score': 'first',
    }).reset_index().sort_values('Performance Score', ascending=False)

    team_avg_prev = prev_month_df['Performance Score'].mean() if len(prev_month_df) > 0 else None
    team_avg_curr = curr_month_df['Performance Score'].mean()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ─── SLIDE 1: COVER ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_BLUE
    top_bar.line.fill.background()

    badge = add_rounded_card(slide, Inches(0.8), Inches(1.5), Inches(2.2), Inches(0.45), ACCENT_BLUE, ACCENT_BLUE)
    tf = badge.text_frame
    tf.paragraphs[0].text = "MARKETING DEPARTMENT"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8), Inches(1), "Performance Analysis", font_size=44, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8), Inches(0.6), f"{period_label} — Executive Summary", font_size=22, color=TEXT_MUTED)

    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.1), Inches(4), Inches(0.04))
    deco.fill.solid()
    deco.fill.fore_color.rgb = ACCENT_BLUE
    deco.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(8), Inches(0.4), "A data-driven narrative of how the Marketing team performed across Digital, Creative, and Content functions", font_size=13, color=TEXT_MUTED)

    add_kpi_card(slide, Inches(9.5), Inches(1.8), Inches(3), Inches(1.3), "TEAM PERFORMANCE", f"{team_avg_curr:.1%}", f"{'↑' if team_avg_prev and team_avg_curr > team_avg_prev else '↓'} May: {team_avg_prev:.1%}" if team_avg_prev else "", ACCENT_GREEN if team_avg_curr >= 0.80 else ACCENT_AMBER)
    add_kpi_card(slide, Inches(9.5), Inches(3.4), Inches(3), Inches(1.3), "TEAM SIZE", f"{curr_month_df['Employee Name'].nunique()} Members", f"Grew from 8 in May", ACCENT_BLUE)
    
    top_performer = perf_curr.iloc[0] if len(perf_curr) > 0 else None
    if top_performer is not None:
        add_kpi_card(slide, Inches(9.5), Inches(5.0), Inches(3), Inches(1.3), "TOP ACHIEVER", f"{top_performer['Employee Name'].split(' ')[0]}", f"Score: {top_performer['Performance Score']:.1%} | {top_performer['Position']}", ACCENT_GREEN)

    # ─── SLIDE 2: THE STORY ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.5), Inches(8), "📖 Executive Summary", f"What happened in {period_label}?", ACCENT_PURPLE)

    story = (
        f"In {period_label}, the Marketing department delivered an average performance of {team_avg_curr:.1%}, "
        f"a slight increase from May's {f'{team_avg_prev:.1%}' if team_avg_prev is not None else 'N/A'}. The team also expanded to 11 members.\n\n"
        f"The standout performer was {top_performer['Employee Name'] if top_performer is not None else 'N/A'} "
        f"with a remarkable {top_performer['Performance Score']:.1%} score, followed closely by Graphic Designers "
        f"who showed massive improvement in project delivery rates compared to May.\n\n"
        f"However, the team faces a deepening crisis in Lead Generation & Media Buying. "
        f"Bahy Hamed Amer's performance plummeted to 51.1% (down from 71.2%), achieving only 35.7% of lead targets "
        f"and 0% of app install targets.\n\n"
        f"On a positive note, the Learning & Growth perspective saw a massive jump to 75% "
        f"(from 35% in May) showing that last month's compliance warnings were taken seriously. "
        f"Social Media Response Time remains a critical bottleneck across all specialists."
    )
    add_rounded_card(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.2))
    add_text_box(slide, Inches(0.9), Inches(1.7), Inches(7), Inches(4.8), story, font_size=14, color=TEXT_DARK)

    add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4), Inches(0.3), "KEY TAKEAWAYS", font_size=12, color=ACCENT_BLUE, bold=True)
    takeaways = [
        ("✅ Design Turnaround", "Graphic designers fixed May volume gaps", ACCENT_GREEN),
        ("⚠️ Response Time", "Still critically missing SLA (4x slower)", ACCENT_RED),
        ("🔴 Media Crisis", "Bahy dropped to 51.1% (Leads & Installs failing)", ACCENT_RED),
        ("📈 L&D Improvement", "Training jumped from 0% to 75% completion", ACCENT_GREEN),
        ("🎯 New Hires", "Abdulwahab joined strong with 91.2%", ACCENT_GREEN),
        ("✅ Quality Maintained", "Brand consistency & creative delivery at 100%", ACCENT_GREEN),
    ]
    for i, (title, desc, color) in enumerate(takeaways):
        y = Inches(2.0) + Inches(0.8) * i
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), y, Inches(0.06), Inches(0.55))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()
        add_text_box(slide, Inches(8.8), y, Inches(4), Inches(0.3), title, font_size=11, color=TEXT_DARK, bold=True)
        add_text_box(slide, Inches(8.8), y + Inches(0.25), Inches(4), Inches(0.3), desc, font_size=10, color=TEXT_MUTED)

    # ─── SLIDE 3: SCOREBOARD CHART ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "🏆 Individual Performance Scoreboard", f"{period_label} — Ranked by Performance Score", ACCENT_CYAN)

    chart_data = CategoryChartData()
    chart_data.categories = [row['Employee Name'].split(' ')[0] + ' ' + row['Employee Name'].split(' ')[-1] for _, row in perf_curr.sort_values('Performance Score', ascending=True).iterrows()]
    chart_data.add_series('Score', [row['Performance Score'] for _, row in perf_curr.sort_values('Performance Score', ascending=True).iterrows()])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.6), Inches(1.5), Inches(8.5), Inches(5.5), chart_data).chart
    chart.has_legend = False
    series = chart.series[0]
    series.has_data_labels = True
    series.data_labels.number_format = '0.0%'
    series.data_labels.font.size = Pt(11)
    series.data_labels.font.color.rgb = TEXT_DARK
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.tick_labels.number_format = '0%'

    dist_left = Inches(9.8)
    add_text_box(slide, dist_left, Inches(1.6), Inches(3), Inches(0.3), "DISTRIBUTION", font_size=12, color=ACCENT_CYAN, bold=True)
    grade_counts = {}
    for _, row in perf_curr.iterrows():
        g, _ = grade_from_score(row['Performance Score'])
        grade_counts[g] = grade_counts.get(g, 0) + 1

    y_dist = Inches(2.1)
    for grade_label in ['A+', 'A', 'B+', 'B', 'C+', 'C', 'D']:
        count = grade_counts.get(grade_label, 0)
        if count > 0:
            _, gc = grade_from_score({'A+': 0.96, 'A': 0.92, 'B+': 0.87, 'B': 0.82, 'C+': 0.77, 'C': 0.72, 'D': 0.5}[grade_label])
            badge_bg = add_rounded_card(slide, dist_left, y_dist, Inches(0.6), Inches(0.35), gc, gc)
            tf = badge_bg.text_frame
            tf.paragraphs[0].text = grade_label
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            add_text_box(slide, dist_left + Inches(0.75), y_dist + Inches(0.05), Inches(2), Inches(0.3),
                         f"{count} {'member' if count == 1 else 'members'}", font_size=11, color=TEXT_DARK)
            y_dist += Inches(0.5)

    add_kpi_card(slide, dist_left, Inches(5.5), Inches(3), Inches(1.3),
                 "TEAM AVERAGE", f"{team_avg_curr:.1%}",
                 f"Target: 80% | Gap: {(team_avg_curr - 0.80)*100:+.1f}%",
                 ACCENT_GREEN if team_avg_curr >= 0.80 else ACCENT_AMBER)

    # ─── SLIDE 4: BSC PERSPECTIVE ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "🎯 Balanced Scorecard — Perspective View", f"How did Marketing perform across each BSC dimension in {period_label}?", ACCENT_GREEN)

    perspectives_data = {
        'Financial': {'avg': 0.7377, 'icon': '💰', 'desc': 'Revenue, CPL, CPV, Budget Compliance', 'insight': 'Declined from May (78%). Media Buyer dropped to 44% in CPL and 35% in CPV causing a huge drag, despite SEO Manager maintaining strong revenue.', 'color': ACCENT_AMBER, 'kpis': [('CPL (Media Buyer)', '44.1%', ACCENT_RED), ('CPV (Media Buyer)', '35.4%', ACCENT_RED), ('Revenue (SEO Manager)', '100%', ACCENT_GREEN)]},
        'Customer': {'avg': 0.8163, 'icon': '👥', 'desc': 'Leads, Volume, App Installs, Campaign Reach', 'insight': 'Stable vs May (82%). Lead generation remains a crisis (35.7%), and app installs dropped to 0% for Media Buying. SEO Manager keeps the average afloat.', 'color': ACCENT_BLUE, 'kpis': [('Leads (Media Buyer)', '35.7%', ACCENT_RED), ('App Installs (Media Buyer)', '0%', ACCENT_RED), ('Campaign Reach', '100%', ACCENT_GREEN), ('Leads (SEO)', '100%', ACCENT_GREEN)]},
        'Internal Process': {'avg': 0.8161, 'icon': '⚙️', 'desc': 'CR, Delivery, Quality, Response Time, Projects', 'insight': 'Slight dip from May (84%). Response Time remains abysmal (<33%). However, Graphic Designers dramatically improved their Project Delivery volume.', 'color': ACCENT_BLUE, 'kpis': [('Brand Consistency', '100%', ACCENT_GREEN), ('Response Time (Social)', '4-33%', ACCENT_RED), ('Edit Rates (Designers)', '35-46%', ACCENT_AMBER)]},
        'Learning & Growth': {'avg': 0.7500, 'icon': '📚', 'desc': 'Training Hours, Research Hours, Development', 'insight': 'Massive turnaround! Jumped from 35% in May to 75% in June. The team responded well to last month\'s L&D warnings.', 'color': ACCENT_GREEN, 'kpis': [('Team Training & Development', '75%', ACCENT_GREEN)]}
    }

    card_width = Inches(5.8)
    card_height = Inches(2.6)
    cards_positions = [(Inches(0.6), Inches(1.5)), (Inches(6.8), Inches(1.5)), (Inches(0.6), Inches(4.4)), (Inches(6.8), Inches(4.4))]

    for idx, (persp_name, data) in enumerate(perspectives_data.items()):
        x, y = cards_positions[idx]
        card = add_rounded_card(slide, x, y, card_width, card_height, CARD_BG)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.3), f"{data['icon']} {persp_name}", font_size=15, color=TEXT_DARK, bold=True)
        add_text_box(slide, x + card_width - Inches(1.5), y + Inches(0.1), Inches(1.3), Inches(0.3), f"{data['avg']:.1%}", font_size=18, color=data['color'], bold=True, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), card_width - Inches(0.4), Inches(0.25), data['desc'], font_size=10, color=TEXT_MUTED)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.75), card_width - Inches(0.4), Inches(0.02))
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER_COLOR
        div.line.fill.background()
        
        for ki, (kpi_name, kpi_val, kpi_color) in enumerate(data['kpis'][:4]):
            ky = y + Inches(0.85) + Inches(0.3) * ki
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), ky + Inches(0.08), Inches(0.1), Inches(0.1))
            dot.fill.solid()
            dot.fill.fore_color.rgb = kpi_color
            dot.line.fill.background()
            
            add_text_box(slide, x + Inches(0.5), ky, Inches(3), Inches(0.25), kpi_name, font_size=10, color=TEXT_DARK)
            add_text_box(slide, x + card_width - Inches(1.2), ky, Inches(1), Inches(0.25), kpi_val, font_size=10, color=kpi_color, bold=True, alignment=PP_ALIGN.RIGHT)
        
        add_text_box(slide, x + Inches(0.2), y + Inches(2.05), card_width - Inches(0.4), Inches(0.6), data['insight'], font_size=9, color=TEXT_MUTED)

    # ─── SLIDE 5: CRITICAL GAPS ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "🔴 Deep Dive — Critical Performance Gaps", f"Root cause analysis for KPIs falling below 70% in {period_label}", ACCENT_RED)

    gaps_data = [
        ("Bahy Hamed", "Media Buyer", "Leads & App Installs", "3,518 / 0", "35% / 0%", "Critical", "Pipeline collapse. App install tracking might be broken or campaigns stopped."),
        ("Bahy Hamed", "Media Buyer", "CPL & CPV", "136 / 401", "44% / 35%", "Critical", "Cost metrics exploded compared to May. Ad spend inefficiency is severe."),
        ("Hossam E.", "Social Head", "Response Time", "89 min", "33.7%", "High", "Slightly better than May but still failing the 30min SLA."),
        ("Samya & Mina", "Social Spc.", "Response Time", "89 min", "4.5%", "Critical", "New hires struggling heavily with response SLA workflows."),
        ("Ahmed Shehata", "SEO Mgr", "Website Speed", "30", "54.5%", "Medium", "Speed score dropped from May (61%), needs technical intervention."),
        ("Nada & Abdul", "Designers", "Edits Rate", "0.20 / 0.15", "35% / 46%", "Medium", "While volume improved, edit rates are now too high indicating rushed work."),
        ("M. Abdelmegid", "Creative Dir.", "Quantity", "100", "66.7%", "Medium", "Same gap as May; personal output volume remains low."),
    ]

    table_top = Inches(1.5)
    col_widths = [Inches(1.5), Inches(1.5), Inches(1.6), Inches(1.0), Inches(1.0), Inches(0.8), Inches(4.5)]
    headers = ["Employee", "Position", "KPI", "Actual", "Ach%", "Severity", "Root Cause Analysis"]
    
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), table_top - Inches(0.05), Inches(12.3), Inches(0.4))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
    header_bg.line.color.rgb = BORDER_COLOR

    x_pos = Inches(0.6)
    for j, (header, width) in enumerate(zip(headers, col_widths)):
        add_text_box(slide, x_pos, table_top, width, Inches(0.3), header, font_size=10, color=TEXT_DARK, bold=True)
        x_pos += width

    for i, gap in enumerate(gaps_data):
        row_y = table_top + Inches(0.45) + Inches(0.65) * i
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), row_y - Inches(0.05), Inches(12.3), Inches(0.6))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
            row_bg.line.fill.background()
        
        x_pos = Inches(0.6)
        for j, (val, width) in enumerate(zip(gap, col_widths)):
            color = TEXT_MUTED
            if j == 4:
                ach_val = float(val.split('%')[0].strip()) if '%' in val else 100
                color = ACCENT_RED if ach_val < 50 else ACCENT_AMBER
            elif j == 5:
                color = ACCENT_RED if val == "Critical" else (ACCENT_AMBER if val == "High" else ACCENT_BLUE)
            elif j == 0:
                color = TEXT_DARK
            add_text_box(slide, x_pos, row_y + Inches(0.05), width, Inches(0.55), val, font_size=10, color=color, bold=(j in (0, 4, 5)))
            x_pos += width

    # ─── SLIDE 6: WINS & EXCELLENCE ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "🌟 Wins & Excellence — What Went Right", f"Spotlighting top performers and 100% KPI achievements in {period_label}", ACCENT_GREEN)

    add_rounded_card(slide, Inches(0.6), Inches(1.5), Inches(4.5), Inches(5.2), RGBColor(240, 253, 244), ACCENT_GREEN)
    add_text_box(slide, Inches(0.9), Inches(1.7), Inches(4), Inches(0.3), "⭐ TOP PERFORMER SPOTLIGHT", font_size=12, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.2), Inches(4), Inches(0.4), "Dina Samir", font_size=24, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.7), Inches(4), Inches(0.3), "Account Manager  |  Score: 94.0%  |  Grade: A+", font_size=12, color=ACCENT_GREEN, bold=True)

    web_dev_kpis = [
        ("Campaign Delivery", "Maintained 100% execution", "100%", ACCENT_GREEN),
        ("Delivery Ontime", "Improved consistency", "94.4%", ACCENT_GREEN),
        ("Modifications", "Zero defect threshold", "100%", ACCENT_GREEN),
        ("Projects Ontime", "Reliable delivery", "90.5%", ACCENT_GREEN),
    ]
    for i, (kpi, detail, ach, color) in enumerate(web_dev_kpis):
        ky = Inches(3.4) + Inches(0.6) * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), ky + Inches(0.05), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text_box(slide, Inches(1.25), ky - Inches(0.03), Inches(2), Inches(0.25), kpi, font_size=11, color=TEXT_DARK, bold=True)
        add_text_box(slide, Inches(1.25), ky + Inches(0.25), Inches(2), Inches(0.2), detail, font_size=10, color=TEXT_MUTED)
        add_text_box(slide, Inches(3.8), ky, Inches(1), Inches(0.25), ach, font_size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

    add_text_box(slide, Inches(5.6), Inches(1.5), Inches(7), Inches(0.3), "NOTABLE 100% ACHIEVEMENTS", font_size=12, color=ACCENT_GREEN, bold=True)
    excellence_items = [
        ("Ahmed Mohamed Shehata", "Leads, Revenue, and App Installs all hit 100% (covering Media Buying gaps)"),
        ("Mohamed Abdelmegid Diab", "Brand Consistency & Creative Delivery maintained perfect scores"),
        ("Hossam El-Naggar", "Campaign Reach & Channels Growth fully achieved again"),
        ("Abdelrahman Yousry", "Website Uptime & Request Delivery remained flawless"),
        ("Samya Ismail", "Response Rate hit 100% despite response time challenges"),
        ("L&D Turnaround", "Department-wide training surged to 75% compliance"),
        ("Graphic Design Output", "Asser & Nada significantly increased volume output vs May"),
    ]
    for i, (name, detail) in enumerate(excellence_items):
        ey = Inches(2.0) + Inches(0.55) * i
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), ey + Inches(0.05), Inches(0.12), Inches(0.12))
        check.fill.solid()
        check.fill.fore_color.rgb = ACCENT_GREEN
        check.line.fill.background()
        add_text_box(slide, Inches(5.85), ey, Inches(2.5), Inches(0.25), name, font_size=11, color=TEXT_DARK, bold=True)
        add_text_box(slide, Inches(5.85), ey + Inches(0.25), Inches(6.5), Inches(0.25), detail, font_size=10, color=TEXT_MUTED)

    # ─── SLIDE 7: MONTH-OVER-MONTH TREND CHART ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "📊 Month-over-Month Trend", f"Comparing {period_label} performance against May baseline", ACCENT_BLUE)

    common_emps = set(may['Employee Name'].unique()).intersection(set(jun['Employee Name'].unique()))
    trend_data = {}
    for emp in sorted(list(common_emps)):
        p_val = may[may['Employee Name'] == emp]['Performance Score'].iloc[0] if len(may[may['Employee Name'] == emp]) > 0 else 0
        c_val = jun[jun['Employee Name'] == emp]['Performance Score'].iloc[0] if len(jun[jun['Employee Name'] == emp]) > 0 else 0
        short_name = emp.split(' ')[0] + ' ' + emp.split(' ')[-1]
        trend_data[short_name] = (p_val, c_val)

    chart_data = CategoryChartData()
    chart_data.categories = list(trend_data.keys())
    chart_data.add_series('May 2026', [d[0] for d in trend_data.values()])
    chart_data.add_series('June 2026', [d[1] for d in trend_data.values()])

    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.4), Inches(12), Inches(4.5), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(12)
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.tick_labels.number_format = '0%'

    add_kpi_card(slide, Inches(0.6), Inches(6.1), Inches(3.5), Inches(1.1), "TEAM AVG (MAY)", f"{team_avg_prev:.1%}", "Baseline", ACCENT_BLUE)
    add_kpi_card(slide, Inches(4.5), Inches(6.1), Inches(3.5), Inches(1.1), "TEAM AVG (JUNE)", f"{team_avg_curr:.1%}", f"{'↑' if team_avg_curr > team_avg_prev else '↓'} {(team_avg_curr - team_avg_prev)*100:+.1f}%", ACCENT_GREEN if team_avg_curr > team_avg_prev else ACCENT_AMBER)
    improved_count = sum(1 for d in trend_data.values() if d[1] > d[0])
    add_kpi_card(slide, Inches(8.4), Inches(6.1), Inches(4.2), Inches(1.1), "IMPROVEMENT RATE", f"{improved_count} / {len(common_emps)}", f"Core members improved in June", ACCENT_GREEN)

    # ─── SLIDE 8: ROLE-BASED ANALYSIS ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "👤 Role-Based Performance Analysis", "How does performance vary by marketing function?", ACCENT_PURPLE)

    roles_analysis = [
        {'role': 'Account Manager', 'member': 'Dina Samir', 'score': 0.9405, 'strengths': 'Consistent delivery, quality control', 'gaps': 'Budget deficit (50%) persists', 'color': ACCENT_GREEN},
        {'role': 'Graphic Designer', 'member': 'Asser Mohamed', 'score': 0.9350, 'strengths': 'Massive volume turnaround vs May', 'gaps': 'Edits rate still needs work', 'color': ACCENT_GREEN},
        {'role': 'Graphic Designer', 'member': 'Abdulwahab', 'score': 0.9120, 'strengths': 'Strong new hire, good volume', 'gaps': 'Edits rate (46%) is very high', 'color': ACCENT_GREEN},
        {'role': 'SEO Manager', 'member': 'Ahmed M. Shehata', 'score': 0.9025, 'strengths': 'Revenue, Leads, App Installs (100%)', 'gaps': 'Website speed dropped to 54%', 'color': ACCENT_GREEN},
        {'role': 'Web Developer', 'member': 'Abdelrahman Y.', 'score': 0.8875, 'strengths': 'Infrastructure reliability (100%)', 'gaps': 'Slight overall dip vs May', 'color': ACCENT_CYAN},
        {'role': 'Creative Director', 'member': 'M. Abdelmegid', 'score': 0.8783, 'strengths': 'Brand consistency, Fixed L&D', 'gaps': 'Personal output volume (66%)', 'color': ACCENT_CYAN},
        {'role': 'Social Media Head', 'member': 'Hossam E.', 'score': 0.8045, 'strengths': 'Reach & growth perfectly met', 'gaps': 'Response time SLA still failing', 'color': ACCENT_BLUE},
        {'role': 'Media Buyer', 'member': 'Bahy Hamed', 'score': 0.5110, 'strengths': 'None in June (Critical drop)', 'gaps': 'Leads (35%), Installs (0%), CPL (44%)', 'color': ACCENT_RED},
    ]

    card_w = Inches(3.0)
    card_h = Inches(2.6)
    for i, role_data in enumerate(roles_analysis):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + (card_w + Inches(0.15)) * col
        y = Inches(1.5) + (card_h + Inches(0.2)) * row
        card = add_rounded_card(slide, x, y, card_w, card_h, CARD_BG)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.25), role_data['role'], font_size=11, color=role_data['color'], bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), card_w - Inches(0.3), Inches(0.2), role_data['member'], font_size=10, color=TEXT_DARK)
        
        grade, gc = grade_from_score(role_data['score'])
        add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(1.5), Inches(0.35), f"{role_data['score']:.1%}", font_size=20, color=TEXT_DARK, bold=True)
        
        grade_b = add_rounded_card(slide, x + Inches(1.8), y + Inches(0.75), Inches(0.5), Inches(0.3), gc, gc)
        tf = grade_b.text_frame
        tf.paragraphs[0].text = grade
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), y + Inches(1.15), card_w - Inches(0.3), Inches(0.015))
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER_COLOR
        div.line.fill.background()
        
        add_text_box(slide, x + Inches(0.15), y + Inches(1.25), card_w - Inches(0.3), Inches(0.15), "STRENGTHS", font_size=9, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.45), card_w - Inches(0.3), Inches(0.4), role_data['strengths'], font_size=9, color=TEXT_MUTED)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.95), card_w - Inches(0.3), Inches(0.15), "GAPS", font_size=9, color=ACCENT_RED, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(2.15), card_w - Inches(0.3), Inches(0.4), role_data['gaps'], font_size=9, color=TEXT_MUTED)

    # ─── SLIDE 9: RECOMMENDATIONS ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(8), "🎯 Strategic Recommendations & Action Plan", "Prioritized interventions to drive Marketing performance in July", ACCENT_BLUE)

    recommendations = [
        {'priority': 'P0', 'title': 'Media Buying Intervention', 'owner': 'Bahy + Management', 'timeline': 'Immediate', 'detail': 'Bahy\'s performance collapsed to 51%. Leads at 35%, CPL exploding. Requires immediate campaign audit and potential reallocation of ad spend.', 'color': ACCENT_RED},
        {'priority': 'P0', 'title': 'Social Response SLA Crisis', 'owner': 'Hossam + New Hires', 'timeline': 'Immediate', 'detail': 'New hires (Samya, Mina) scoring ~4% on response time SLA. Process is clearly broken. Needs automation or shift-scheduling redesign.', 'color': ACCENT_RED},
        {'priority': 'P1', 'title': 'Quality Control for Design', 'owner': 'M. Abdelmegid', 'timeline': '15 Days', 'detail': 'While designers fixed their volume issues from May, their Edit Rates are now failing (35-46%). Need to balance speed with first-time accuracy.', 'color': ACCENT_AMBER},
        {'priority': 'P1', 'title': 'Website Speed Drop', 'owner': 'Abdelrahman + Ahmed', 'timeline': '30 Days', 'detail': 'Speed score dropped to 54.5%. Needs technical audit to ensure recent web updates aren\'t harming SEO metrics.', 'color': ACCENT_AMBER},
        {'priority': 'P2', 'title': 'Budget Deficit Management', 'owner': 'Dina Samir', 'timeline': '30 Days', 'detail': 'Deficit remains at 50% achievement for the second month. Need tighter cost control mechanisms on active campaigns.', 'color': ACCENT_BLUE},
    ]

    for i, rec in enumerate(recommendations):
        y = Inches(1.4) + Inches(1.05) * i
        p_badge = add_rounded_card(slide, Inches(0.6), y + Inches(0.05), Inches(0.6), Inches(0.35), rec['color'], rec['color'])
        tf = p_badge.text_frame
        tf.paragraphs[0].text = rec['priority']
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_text_box(slide, Inches(1.4), y, Inches(5), Inches(0.3), rec['title'], font_size=14, color=TEXT_DARK, bold=True)
        meta = f"Owner: {rec['owner']}  |  Timeline: {rec['timeline']}"
        add_text_box(slide, Inches(1.4), y + Inches(0.3), Inches(6), Inches(0.2), meta, font_size=10, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, Inches(1.4), y + Inches(0.55), Inches(11), Inches(0.3), rec['detail'], font_size=11, color=TEXT_MUTED)

    # ─── SLIDE 10: BOTTOM LINE ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_divider(slide, Inches(0.8), Inches(0.5), Inches(8), "📌 The Bottom Line", f"{period_label} — Marketing Department Performance Summary", ACCENT_PURPLE)

    add_rounded_card(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(3.2), CARD_BG)
    closing_narrative = (
        f"The Marketing team in {period_label} delivered an overall performance of 82.2% — a slight "
        "improvement from May, fueled by strong L&D adherence and Design team volume turnaround.\n\n"
        "THE QUALITY & GROWTH STORY: Brand consistency is flawless. Training compliance jumped from 0% to 75%. "
        "SEO Manager stepped up to cover revenue gaps.\n\n"
        "THE CRISIS POINTS: Media Buying (Bahy) has completely collapsed regarding leads and cost efficiency. "
        "Social Media Response Time is broken, with new hires severely missing SLAs (4%). "
        "Design volume improved, but at the cost of higher edit/rework rates.\n\n"
        "ACTION: July requires urgent surgical intervention on Media Buying campaigns and a complete "
        "overhaul of how the Social Media team handles incoming response queues."
    )
    add_text_box(slide, Inches(1.8), Inches(1.9), Inches(9.4), Inches(2.8), closing_narrative, font_size=13, color=TEXT_DARK)

    kpi_cards_data = [
        ("INTERNAL & L&D", "81%", "Process & Training", ACCENT_GREEN),
        ("CUSTOMER", "81.6%", "Saved by SEO", ACCENT_BLUE),
        ("FINANCIAL", "73.7%", "Dragged by Ads", ACCENT_RED),
        ("OVERALL", "82.2%", "Team Average", ACCENT_BLUE),
    ]
    for i, (label, value, desc, color) in enumerate(kpi_cards_data):
        x = Inches(0.8) + Inches(3.1) * i
        add_kpi_card(slide, x, Inches(5.2), Inches(2.8), Inches(1.5), label, value, desc, color)

    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                 f"Marketing PMS Analysis  •  {period_label}  •  Confidential Corporate Report",
                 font_size=10, color=TEXT_LIGHT)

    # ─── SAVE TO BYTES ───
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
