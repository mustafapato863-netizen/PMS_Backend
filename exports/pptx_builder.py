import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def build_pptx_from_slides(report_name: str, slides_data: list, period_label: str) -> bytes:
    """
    Builds a PowerPoint presentation using python-pptx based on structured slide data.
    """
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = report_name
    subtitle.text = f"Monthly Performance Report\n{period_label}"

    # Standard Blank layout for content slides
    blank_layout = prs.slide_layouts[5] # Title Only layout

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Set slide title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

        blocks = slide_data.get("blocks", [])

        # Very simple deterministic layout engine
        # For now, we will vertically stack the blocks
        current_top = Inches(2.0)
        left_margin = Inches(0.5)
        width = Inches(9.0)

        for block in blocks:
            b_type = block.get("type")
            settings = block.get("config", {}).get("settings", {})
            block_title = settings.get("title", "")

            if b_type == "narrative":
                txBox = slide.shapes.add_textbox(left_margin, current_top, width, Inches(1.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = block_title or "Narrative content will be generated here."
                p.font.size = Pt(14)
                current_top += Inches(1.8)

            elif b_type == "kpi_summary":
                # Render 3 boxes side by side
                box_width = Inches(2.8)
                for i in range(3):
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        left_margin + (Inches(3.0) * i),
                        current_top,
                        box_width,
                        Inches(1.2)
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(240, 244, 248) # Slate-50 approx
                    shape.line.color.rgb = RGBColor(203, 213, 225) # Slate-200

                    tf = shape.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"KPI {i+1}\n85%"
                    p.alignment = PP_ALIGN.CENTER
                    p.font.color.rgb = RGBColor(15, 23, 42) # Slate-900

                current_top += Inches(1.5)

            elif b_type == "data_table":
                rows = 4
                cols = 2
                table_shape = slide.shapes.add_table(rows, cols, left_margin, current_top, width, Inches(1.5))
                table = table_shape.table
                table.columns[0].width = Inches(4.5)
                table.columns[1].width = Inches(4.5)

                # Header
                table.cell(0, 0).text = "Metric"
                table.cell(0, 1).text = "Value"

                for i in range(1, 4):
                    table.cell(i, 0).text = f"Sample Data {i}"
                    table.cell(i, 1).text = "100"

                current_top += Inches(2.0)

            elif b_type in ("bar_chart", "line_chart"):
                # Placeholder for chart since we don't have real chart data attached to the block yet
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left_margin,
                    current_top,
                    width,
                    Inches(3.0)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
                shape.line.color.rgb = RGBColor(226, 232, 240)

                tf = shape.text_frame
                p = tf.paragraphs[0]
                p.text = f"[{b_type.replace('_', ' ').title()}: {block_title}]"
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor(148, 163, 184)

                current_top += Inches(3.2)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
